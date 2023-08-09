import os
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import pandas as pd
import base64
import io
from getter import connexion


def get_wafer(wafer_id):
    """
    This function finds the wafer specified in the database

    :param <str> wafer_id: name of the wafer_id

    :return <dict>: the wafer
    """
    collection = connexion()
    return collection.find_one({"wafer_id": wafer_id})


def plot_wanted_matrices(wafer_id, sessions, structures, types, Temps, Files, coords):
    """
    Used to plot only selected dies following selected filters. These plots won't be registered but will be displayed for the user.
    One plot is created per type of Measurement.
    A single plot contain all session, all dies and all structures selected. A list of 15 colors is generated, so up to 15 differents lines can be differenciated.
    You can add a color in this list if you want (Line 44 to 60) but don't forget to change the number in line 119 (color_index % <this number>)

    :param wafer_id: ID of the wafer
    :param sessions: Selected sessions
    :param structures: Selected structures
    :param types: Selected types
    :param Temps: Selected temperatures
    :param Files: Selected files
    :param coords: Selected coordinates

    :return: One plot per type of measurements, converted to base64
    """
    wafer = get_wafer(wafer_id)
    figs = []
    if 'C' in types:
        types.remove('C')
    colors = [
        '#FF0000',  # Rouge
        '#00FF00',  # Vert
        '#0000FF',  # Bleu
        '#FFFF00',  # Jaune
        '#FF00FF',  # Magenta
        '#00FFFF',  # Cyan
        '#800000',  # Marron
        '#808000',  # Olive
        '#008000',  # Vert foncé
        '#008080',  # Vert d'eau
        '#000080',  # Bleu marine
        '#800080',  # Violet
        '#7F7F7F',  # Gris
        '#FF6600',  # Orange
        '#663399'  # Rebecca Purple
    ]

    y_values_dict = {'I': ['I'], 'J': ['J'], 'It': ['It']}

    df_dict = {}
    for type in types:
        df_dict[type] = pd.DataFrame()
        y_values = y_values_dict[type]
        for session in sessions:
            for structure in structures:
                if wafer[session].get(structure) is None or structure == 'Compliance':
                    continue
                for matrix in wafer[session][structure]["matrices"]:
                    coordinates = f"({matrix['coordinates']['x']},{matrix['coordinates']['y']})"
                    if coordinates in coords:
                        results = [result for result in matrix["results"]]
                        if type in results:
                            if matrix["results"][type]["Filename"] in Files and matrix["results"][type]["Temperature"] in Temps:
                                if type != 'It':
                                    data_values = [[] for _ in range(len(y_values) + 1)]
                                    data_values[0] = ['Voltage (V)']

                                    for double in matrix["results"][type]["Values"]:
                                        data_values[0].append(abs(double["V"]))
                                        for idx, unit in enumerate(y_values, 1):
                                            if double.get(unit) is not None:
                                                data_values[idx].append(abs(double[unit]))

                                else:
                                    data_values = [[] for _ in range(len(y_values) + 1)]
                                    data_values[0] = ['T (seconds)']

                                    for double in matrix["results"][type]["Values"]:
                                        data_values[0].append(abs(double["t (s)"]))
                                        for idx, unit in enumerate(y_values, 1):
                                            if double.get(unit) is not None:
                                                data_values[idx].append(abs(double[unit]))

                                columns = [f"{' '.join(session.split(' ')[1:])} {coordinates} {structure} {matrix['results'][type]['Temperature']} {matrix['results'][type]['Filename']}"] + [f"{' '.join(session.split(' ')[1:])} {coordinates} {structure} {matrix['results'][type]['Temperature']} {matrix['results'][type]['Filename']}" + " " + unit for unit in y_values]
                                new_df = pd.DataFrame(list(zip(*data_values)), columns=columns)
                                df_dict[type] = df_dict[type].merge(new_df, left_index=True, right_index=True, how='outer')

    for unit, df in df_dict.items():
        cols = [col.rstrip() for col in df.columns if not col.endswith(' ' + unit)]
        if not cols:  # If no columns for this unit, skip
            continue

        fig, ax = plt.subplots(figsize=(10, 5))  # Crée une figure avec une taille spécifique
        color_index = 0

        for col in cols:
            if not df[col].empty:
                label_x = df[col].iloc[0]
                label_y = f"{unit} Value ({unit})"

                df[col] = df[col].iloc[1:]
                df[col + ' ' + unit] = df[col + ' ' + unit].iloc[1:]
                df[col + ' ' + unit] = abs(df[col + ' ' + unit])

                ax.plot(df[col], df[col + ' ' + unit], str(colors[color_index % 15]), label=col)

                ax.set_xlabel(label_x)
                ax.set_ylabel(label_y)

                plt.yscale('log')

                ax.set_title(f"{wafer_id} {unit}")
                ax.legend(loc='upper left', bbox_to_anchor=(1, 1), fontsize='small')
                ax.grid(True)

                color_index += 1

        plt.subplots_adjust(right=0.7)  # Ajuste la zone de dessin pour laisser de l'espace à la légende
        figs.append(fig)

    figures = []
    for fig in figs:
        figures.append(fig_to_base64(fig))
        plt.close(fig)

    return figures


def wanted_ppt(wafer_id, sessions, structures, types, Temps, Files, coords, file_name):
    """
    Used to create a PowerPoint (registered in a folder named after the concerned wafer) with only selected dies following selected filters. These plots will be registered in a folder named 'plots'.
    One plot is created per type of Measurement.
    A single plot contain all session, all dies and all structures selected. A list of 15 colors is generated, so up to 15 differents lines can be differenciated on the plot.
    You can add a color in this list if you want (Line 169 to 185) but don't forget to change the number in line 245 (color_index % <this number>)

    :param wafer_id: ID of the wafer
    :param sessions: Selected sessions
    :param structures: Selected structures
    :param types: Selected types
    :param Temps: Selected temperatures
    :param Files: Selected files
    :param coords: Selected coordinates
    :param file_name: Name given to the file

    :return: One plot per type of measurements, converted to base64
    """
    if not os.path.exists("plots"):
        os.makedirs("plots")
    if not os.path.exists(wafer_id):
        os.makedirs(wafer_id)

    wafer = get_wafer(wafer_id)
    if 'C' in types:
        types.remove('C')
    colors = [
        '#FF0000',  # Rouge
        '#00FF00',  # Vert
        '#0000FF',  # Bleu
        '#FFFF00',  # Jaune
        '#FF00FF',  # Magenta
        '#00FFFF',  # Cyan
        '#800000',  # Marron
        '#808000',  # Olive
        '#008000',  # Vert foncé
        '#008080',  # Vert d'eau
        '#000080',  # Bleu marine
        '#800080',  # Violet
        '#7F7F7F',  # Gris
        '#FF6600',  # Orange
        '#663399'  # Rebecca Purple
    ]

    y_values_dict = {'I': ['I'], 'J': ['J'], 'It': ['It']}

    df_dict = {}
    for type in types:
        df_dict[type] = pd.DataFrame()
        y_values = y_values_dict[type]
        for session in sessions:
            for structure in structures:
                if wafer[session].get(structure) is None:
                    continue
                for matrix in wafer[session][structure]["matrices"]:
                    coordinates = f"({matrix['coordinates']['x']},{matrix['coordinates']['y']})"
                    if coordinates in coords:
                        results = [result for result in matrix["results"]]
                        if type in results:
                            if matrix["results"][type]["Filename"] in Files and matrix["results"][type]["Temperature"] in Temps:
                                if type != 'It':
                                    data_values = [[] for _ in range(len(y_values) + 1)]
                                    data_values[0] = ['Voltage (V)']

                                    for double in matrix["results"][type]["Values"]:
                                        data_values[0].append(abs(double["V"]))
                                        for idx, unit in enumerate(y_values, 1):
                                            data_values[idx].append(abs(double[unit]))

                                else:
                                    data_values = [[] for _ in range(len(y_values) + 1)]
                                    data_values[0] = ['T (seconds)']

                                    for double in matrix["results"][type]["Values"]:
                                        data_values[0].append(abs(double["t (s)"]))
                                        for idx, unit in enumerate(y_values, 1):
                                            data_values[idx].append(abs(double[unit]))

                                columns = [f"{' '.join(session.split(' ')[1:])} {coordinates} {structure} {matrix['results'][type]['Temperature']} {matrix['results'][type]['Filename']}"] + [f"{' '.join(session.split(' ')[1:])} {coordinates} {structure} {matrix['results'][type]['Temperature']} {matrix['results'][type]['Filename']}" + " " + unit for unit in y_values]
                                new_df = pd.DataFrame(list(zip(*data_values)), columns=columns)
                                df_dict[type] = df_dict[type].merge(new_df, left_index=True, right_index=True, how='outer')

    prs = Presentation()
    units = {"I": "A", "J": "A / cm^2", "It": "s"}

    for unit, df in df_dict.items():
        cols = [col.rstrip() for col in df.columns if not col.endswith(' ' + unit)]
        if not cols:  # If no columns for this unit, skip
            continue

        fig, ax = plt.subplots(figsize=(10, 5))  # Crée une figure avec une taille spécifique
        color_index = 0

        for col in cols:
            if not df[col].empty:
                label_x = df[col].iloc[0]
                label_y = f"{unit} Value ({units[unit]})"

                df[col] = df[col].iloc[1:]
                df[col + ' ' + unit] = df[col + ' ' + unit].iloc[1:]
                df[col + ' ' + unit] = abs(df[col + ' ' + unit])

                ax.plot(df[col], df[col + ' ' + unit], str(colors[color_index % 15]), label=col)

                ax.set_xlabel(label_x)
                ax.set_ylabel(label_y)

                plt.yscale('log')

                ax.set_title(f"{wafer_id} {unit}")
                ax.legend(loc='upper left', bbox_to_anchor=(1, 1), fontsize='small')
                ax.grid(True)

                color_index += 1

        plt.subplots_adjust(right=0.7)  # Ajuste la zone de dessin pour laisser de l'espace à la légende
        plt.savefig(f"plots\\{file_name}_{wafer_id}_{unit}.png")
        plt.close()

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(1)
        top = Inches(1)

        slide.shapes.add_picture(f"plots\\{file_name}_{wafer_id}_{unit}.png", left, top)

    if len(prs.slides) > 0:
        prs.save(f"{wafer_id}\\{file_name}.pptx")
    return


def fig_to_base64(fig):
    """
    Function used to convert a png to base64 to help communication between server and User. Used in ppt_matrix.

    :param <png> fig: a figure in png format

    :return <base64>: The converted figure
    """
    fig_file = io.BytesIO()
    fig.savefig(fig_file, format='png')
    fig_file.seek(0)
    fig_png_base64 = base64.b64encode(fig_file.read())
    fig_file.close()
    return fig_png_base64.decode('utf-8')
